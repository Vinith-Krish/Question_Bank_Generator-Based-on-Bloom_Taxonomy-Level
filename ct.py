import fitz  # PyMuPDF for PDFs
import docx  # python-docx for Word files
from pptx import Presentation  # python-pptx for PowerPoint files
import streamlit as st
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
import google.generativeai as palm
from langchain_google_genai import ChatGoogleGenerativeAI
import PyPDF2
import io
import pandas as pd
import json
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Table, TableStyle
from reportlab.lib.units import inch
 

# Set up Gemini API key
google_api_key = "AIzaSyDUgOcQy1AokHighaXjxFlEYqpu6cSmP9I"
 
st.title("File to Question Bank Generator")
 
# File upload for multiple types
uploaded_file = st.file_uploader("Upload a file (PDF, DOCX, PPTX)", type=["pdf", "docx", "pptx"])
def extract_text_from_pdf(uploaded_file):
    if uploaded_file is not None:
        with fitz.open(stream=uploaded_file.read(), filetype="pdf") as doc:
            text = ""
            for page in doc:
                text += page.get_text()
        return text
    return None
 
 
 
 
# Function to extract text from DOCX (Word)
def extract_text_from_docx(uploaded_file):
    if uploaded_file is not None:
        try:
            doc = docx.Document(uploaded_file)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text
        except Exception as e:
            st.error(f"Error reading DOCX file: {e}")
            return None
    return None
 
# Function to extract text from PPTX (PowerPoint)
def extract_text_from_pptx(uploaded_file):
    if uploaded_file is not None:
        try:
            presentation = Presentation(uploaded_file)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            st.error(f"Error reading PPTX file: {e}")
            return None
    return None
 
# Extract text based on file type
def extract_text(uploaded_file):
    if uploaded_file is not None:
        file_type = uploaded_file.name.split(".")[-1].lower()
        if file_type == "pdf":
            return extract_text_from_pdf(uploaded_file)
        elif file_type == "docx":
            return extract_text_from_docx(uploaded_file)
        elif file_type == "pptx":
            return extract_text_from_pptx(uploaded_file)
        else:
            st.error("Unsupported file type!")
            return None
    return None

taxonomy_level_list= ['Knowledge', 'Comprehension', 'Application', 'Analysis', 'Synthesis', 'Evaluation']
# Dropdown to select Bloom's Taxonomy level
taxonomy_level = st.selectbox("Select Bloom's Taxonomy Level",taxonomy_level_list)

print(taxonomy_level)

final_taxonomy_level = []

for level in taxonomy_level_list:
    final_taxonomy_level.append(level)
    if level == taxonomy_level:
        break

print(final_taxonomy_level)

 
# Course outcome and module input fields
course_outcome = st.text_input("Enter Course Outcome")
module = st.text_input("Enter Modul and Subject Name")
 
# Extracted text from file
extracted_text = extract_text(uploaded_file)
 
# Google Palm API setup for question generation
def get_question_generation_chain():
    prompt_template = """
        Generate a question bank using Bloom's Taxonomy from the provided content.
 
        You need to create 20 questions for each level of Bloom's Taxonomy based on the specified action verbs. Each question should clearly reflect its targeted taxonomy level and adhere to the following specifications.
 
        *Knowledge-Based Questions: Use verbs such as **define, **memorize, **repeat, **copy, **identify, **state, **list, **quote, or **find*.
        *Comprehension-Based Questions: Employ verbs like **summarize, **compare, **describe, **explain, **discuss, **recognize, **report, **translate, and **categorize*.
        *Application-Based Questions: Implement verbs such as **determine, **present, **examine, **implement, **solve, **use, **demonstrate, **interpret, and **reenact*.
        *Analysis-Based Questions: Include verbs such as **organize, **compare, **contrast, **experiment, **test, **question, **connect, **deduce, and **link*.
        *Synthesis-Based Questions: Use verbs like **design, **compose, **construct, **develop, **formulate, **build, **write, or **simulate*.
        *Evaluation-Based Questions: Apply verbs such as **argue, **defend, **judge, **support, **value, **weigh, **reflect, **review, and **grade*.
 
        # Steps
 
        1. Read the provided content {context}.
        2. For each taxonomy level, generate four unique questions employing the related verbs.
        3. Ensure that each question incorporates at least one action verb, which should be in bold.
        4. After composing the question, append its taxonomy level in square brackets at the end.
        5. Each question must conclude with a period.
 
        # Output Format
 
        Produce questions in the following format: "[Question Text] [Taxonomy Level in square brackets with action verbs in *bold*]."
 
        # Examples
 
        ### Knowledge-Based Questions Example:
        What is the primary function of the *heart*? [Knowledge]
        *Identify* the components of the *solar* system. [Knowledge]
 
        ### Comprehension-Based Questions Example:
        *Explain* how photosynthesis occurs in *plants*. [Comprehension]
        *Describe* the *process* of cellular respiration. [Comprehension]
 
        ### Application-Based Questions Example:
        *Solve* for the variable *x* in the equation 2x + 3 = 7. [Application]
        *Demonstrate* how to *prepare* a simple salad. [Application]
 
        ### Analysis-Based Questions Example:
        *Contrast* the *nervous* system with the *endocrine* system. [Analysis]
        *Test* the *effects* of sunlight on plant growth. [Analysis]
 
        ### Synthesis-Based Questions Example:
        *Design* an *experiment* to test water purity. [Synthesis]
        *Compose* a *short* story based on the given theme. [Synthesis]
 
        ### Evaluation-Based Questions Example:
        *Defend* the *importance* of renewable energy. [Evaluation]
        *Judge* the *effectiveness* of the new law on pollution control. [Evaluation]
 
        (Note: For real examples, ensure that questions align with the specific content and context provided.)
    """
   
    model = ChatGoogleGenerativeAI(
        model="gemini-1.5-flash",
        client=palm,
        temperature=0,
        google_api_key=google_api_key
    )
    prompt = PromptTemplate(template=prompt_template, input_variables=["context"])
    chain = LLMChain(llm=model, prompt=prompt)
    return chain
 
# Function to process file and generate questions
def process_file_for_questions(extracted_text):
    if not extracted_text:
        st.error("No text extracted from the file.")
        return None
 
    chain = get_question_generation_chain()
    response = chain.run(context=extracted_text)

    return response
 
# Filtering function
def get_filtering_chain():
    filter_prompt_template = """
       Task:

        You will receive a list of questions categorized under Bloom's Taxonomy, along with selected taxonomy levels. 
        Your objective is to filter and return exactly 20 questions based on the specified levels.

        Instructions:

        Input: A list of questions are in {response} with their corresponding Bloom's Taxonomy levels and 
        the selected levels indicated by {final_taxonomy_level}.

        Filtering:

        Select a total of 20 questions that include all levels specified in {final_taxonomy_level}.
        Do not exceed 20 questions, and ensure there are no duplicates.
        Output Format: Each question should be structured in JSON format, including the question and its corresponding level. For example:

        [ {{ "question": "What is the capital of France?", "level": "Knowledge" }}, {{ "question": "Explain the process of photosynthesis.", "level": "Comprehension" }} ]

        Strictly follow the output format do not add ```json in begining and ``` at the end

        Important: The questions must be ordered starting from Knowledge, followed by Comprehension, Application, Analysis, Synthesis, and Evaluation. Only include levels mentioned in {final_taxonomy_level}; all others should be excluded.   
    """
   
    model = ChatGoogleGenerativeAI(
        model="gemini-1.5-flash",
        client=palm,
        temperature=0,
        google_api_key=google_api_key,
        convert_system_message_to_human=True
    )
    prompt = PromptTemplate(template=filter_prompt_template, input_variables=["response", "final_taxonomy_level"])
    filter_chain = LLMChain(llm=model, prompt=prompt)
    return filter_chain

# New function to save DataFrame as a properly formatted PDF
def save_dataframe_to_pdf(df, module_name, filename):
    doc = SimpleDocTemplate(filename, pagesize=A4, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    styles = getSampleStyleSheet()
    elements = []

    # Add title with module name
    title = Paragraph(f'Question Bank for Module: {module_name}', styles['Title'])
    elements.append(title)

    # Add a line break
    elements.append(Paragraph('<br/>', styles['Normal']))

    # Prepare data for the table
    table_data = [df.columns.tolist()] + [[Paragraph(str(item), styles['Normal']) for item in row] for row in df.values.tolist()]

    # Create a table with adjusted column widths
    table = Table(table_data, colWidths=[1 * inch, 3.5 * inch, 2.5 * inch])  # Adjusted for question numbers

    # Style the table
    table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),  # Header background
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, 0), 'CENTER'),  # Header center alignment
        ('ALIGN', (0, 1), (-1, -1), 'LEFT'),   # Left align all rows
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('TOPPADDING', (0, 1), (-1, -1), 10),  # Add top padding for body rows
        ('BOTTOMPADDING', (0, 1), (-1, -1), 10),  # Add bottom padding for body rows
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),  # Body background
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
    ]))

    # Add table to elements
    elements.append(table)

    # Build the PDF
    doc.build(elements)

# Process file and generate questions when "Generate Questions" button is clicked
if st.button("Generate Questions"):
    with st.spinner("Generating questions..."):
        if extracted_text and final_taxonomy_level:
            # Initialize progress bar
            progress = st.progress(0)

            questions = process_file_for_questions(extracted_text)
            if questions:
                # Update progress
                progress.progress(50)  # Update to 50% after generating questions

                # Filter questions based on selected taxonomy level
                filter_chain = get_filtering_chain()
                filtered_questions = filter_chain.run(response=questions, final_taxonomy_level=final_taxonomy_level)

                # Update progress
                progress.progress(100)

                print(f"Filtered Questions: {filtered_questions}")

                dictionaries = json.loads(filtered_questions)
                df = pd.DataFrame(dictionaries)
                # Add question numbers
                df['Question #'] = range(1, len(df) + 1)

                # Reorder columns to display question numbers first
                df = df[['Question #', 'question', 'level']]

                # Display as table
                st.dataframe(df)

                # Save the DataFrame to a PDF when the button is clicked
                pdf_filename = "questions.pdf"
                save_dataframe_to_pdf(df, module, pdf_filename)

                # Provide a download link for the generated PDF
                with open(pdf_filename, "rb") as pdf_file:
                    st.download_button(
                        label="Download Questions as PDF",
                        data=pdf_file,
                        file_name=pdf_filename,
                        mime="application/pdf"
                    )

            else:
                st.error("Failed to generate questions.")
        else:
            st.error("Please upload a file and select a taxonomy level.")